#!/usr/bin/env python3
"""Generate ShipShape Final Presentation PowerPoint (32 slides)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import os

# --- Colors ---
BG = RGBColor(0x0D, 0x0D, 0x0D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x00, 0x7A, 0xCC)
GREEN = RGBColor(0x4E, 0xC9, 0xB0)
YELLOW = RGBColor(0xDC, 0xDC, 0xAA)
GRAY = RGBColor(0x99, 0x99, 0x99)
RED = RGBColor(0xF4, 0x4E, 0x4E)
TBL_HEADER = RGBColor(0x00, 0x56, 0x8F)
TBL_EVEN = RGBColor(0x1A, 0x1A, 0x2E)
TBL_ODD = RGBColor(0x12, 0x12, 0x20)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK_LAYOUT = prs.slide_layouts[6]  # blank


def set_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG


def tb(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf


def ap(tf, text, size=18, color=WHITE, bold=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_before = Pt(6)
    return p


def _set_cell_bg(cell, rgb_color):
    tcPr = cell._tc.get_or_add_tcPr()
    solidFill = tcPr.makeelement(qn('a:solidFill'), {})
    srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': '%02X%02X%02X' % (rgb_color[0], rgb_color[1], rgb_color[2])})
    solidFill.append(srgbClr)
    tcPr.append(solidFill)


def tbl(slide, left, top, width, height, data, col_widths=None):
    rows, cols = len(data), len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table

    # Remove default banding
    tbl_pr = table._tbl.tblPr
    tbl_pr.attrib['bandRow'] = '0'
    tbl_pr.attrib['bandCol'] = '0'
    tbl_pr.attrib['firstRow'] = '0'
    tbl_pr.attrib['lastRow'] = '0'

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for r_idx, row_data in enumerate(data):
        for c_idx, cell_text in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = WHITE
                if r_idx == 0:
                    paragraph.font.bold = True

            # Background colors
            if r_idx == 0:
                _set_cell_bg(cell, (0x00, 0x56, 0x8F))
            elif r_idx % 2 == 0:
                _set_cell_bg(cell, (0x1A, 0x1A, 0x2E))
            else:
                _set_cell_bg(cell, (0x12, 0x12, 0x20))

            # Remove cell borders
            tcPr = cell._tc.get_or_add_tcPr()
            for border_name in ['lnL', 'lnR', 'lnT', 'lnB']:
                ln = tcPr.makeelement(qn('a:' + border_name), {})
                noFill = ln.makeelement(qn('a:noFill'), {})
                ln.append(noFill)
                tcPr.append(ln)

    return table


# ============================================================
# SLIDE 1: Title
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s)
tb(s, 0.5, 1.5, 12.333, 1.0, "ShipShape", 48, ACCENT, True, PP_ALIGN.CENTER)
tb(s, 0.5, 2.7, 12.333, 0.8, "Performance Audit & Optimization \u2014 Final Report", 28, WHITE, False, PP_ALIGN.CENTER)
tb(s, 0.5, 3.7, 12.333, 0.6, "7 Categories | 24 Specs | All Targets Met", 20, GREEN, False, PP_ALIGN.CENTER)
tb(s, 0.5, 4.5, 12.333, 0.6, "Faheem Syed | Phase 2 + Phase 3 | March 2026", 16, GRAY, False, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: Executive Summary
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s)
tb(s, 0.5, 0.3, 12, 0.7, "Executive Summary", 36, ACCENT, True)
tbl(s, 0.5, 1.2, 12, 5.5, [
    ["Metric", "Result"],
    ["Categories improved", "7 of 7"],
    ["Targets met", "7 of 7"],
    ["Specs implemented", "24 (Phase 2) + 30+ commits (Phase 3)"],
    ["Tests restored", "+1,012 (30% to 99.5%)"],
    ["Type violations", "-44% explicit (-82% any)"],
    ["API response", "-50% wiki p99"],
    ["DB queries", "-36% main page"],
    ["WCAG violations", "13 of 13 fixed"],
    ["Silent failures", "5 of 6 fixed"],
    ["CI pipeline", "8-shard ~5 min"],
], col_widths=[4, 8])

# ============================================================
# SLIDE 3: Methodology
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s)
tb(s, 0.5, 0.3, 12, 0.7, "Methodology", 36, ACCENT, True)
tf = tb(s, 0.5, 1.3, 12, 5.5, "Baseline: Original master branch, seeded DB (501 docs, 218 issues, 22 users, 35 sprints)", 14)
ap(tf, "Branches: Each category on its own branch (cat1 through cat7)", 14)
ap(tf, "Measurement: Before/after under identical conditions, 3 runs median", 14)
ap(tf, "Commits: Structured format \u2014 problem, fix, tradeoffs, measured improvement", 14)
ap(tf, "Merges: --no-ff to master, preserving branch history", 14)
ap(tf, "CI: 8-shard GitHub Actions for reproducible E2E results", 14)

# ============================================================
# SLIDE 4: All 7 Targets Met
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s)
tb(s, 0.5, 0.3, 12, 0.7, "All 7 Targets Met", 36, ACCENT, True)
tbl(s, 0.3, 1.2, 12.7, 5.5, [
    ["Category", "Baseline", "After", "Change", "Target"],
    ["Type Safety", "708 explicit", "397", "-44%", "MET"],
    ["Bundle Size", "2,073 KB monolith", "955 KB largest", "-54%", "MET"],
    ["API Response", "142ms wiki p99", "71ms", "-50%", "MET"],
    ["DB Queries", "25 queries/page", "~16", "-36%", "MET"],
    ["Test Coverage", "451/1,479 (30%)", "1,463/1,471", "+1,012", "EXCEEDED"],
    ["Error Handling", "6 silent failures", "1 remaining", "5 fixed", "EXCEEDED"],
    ["Accessibility", "3 Crit + 10 Serious", "0 + 0", "-13", "MET"],
], col_widths=[2.2, 3.0, 2.5, 2.0, 2.0])

# ============================================================
# SLIDE 5: Cat 1 Before/After
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s)
tb(s, 0.5, 0.3, 12, 0.7, "Cat 1: Type Safety \u2014 MET (44%)", 36, ACCENT, True)
tbl(s, 0.5, 1.2, 12, 4.0, [
    ["Violation", "Before", "After", "Change"],
    ["Explicit any", "392", "70", "-82%"],
    ["Type assertions", "280", "283", "+1%"],
    ["Non-null (!)", "35", "43", "+23%"],
    ["Total explicit", "708", "397", "-44%"],
], col_widths=[4, 2.5, 2.5, 3])

# ============================================================
# SLIDE 6: Cat 1 What Changed
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s)
tb(s, 0.5, 0.3, 12, 0.7, "Cat 1: What Changed", 36, ACCENT, True)
tf = tb(s, 0.5, 1.2, 12, 6.0, "Spec 1.1 \u2014 DB Row Types", 16, GREEN, True)
ap(tf, "Added TypeScript interfaces for all query results in projects.ts and weeks.ts. No more any on database rows.", 14)
ap(tf, "Spec 1.2 \u2014 Discriminated Union", 16, GREEN, True)
ap(tf, "Created union type for 5 document types. TypeScript narrows automatically on switch(doc.type).", 14)
ap(tf, "Spec 1.3 \u2014 Yjs Converter", 16, GREEN, True)
ap(tf, "Added TipTap JSON types to Yjs pipeline. 15 any types eliminated from yjsConverter.ts.", 14)
ap(tf, "Spec 1.4 \u2014 Strict tsconfig", 16, GREEN, True)
ap(tf, "Enabled noUncheckedIndexedAccess and noImplicitReturns. Fixed 102 new errors.", 14)
ap(tf, "Tradeoff: 25 non-null assertions added for bounds-checked array access.", 14, YELLOW)

# ============================================================
# SLIDE 7: Cat 1 How to Reproduce
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s)
tb(s, 0.5, 0.3, 12, 0.7, "Cat 1: How to Reproduce", 36, ACCENT, True)
tf = tb(s, 0.5, 1.2, 12, 5.5, "Count explicit any:", 14, GRAY)
ap(tf, "grep -rn ': any' api/src web/src shared/src --include='*.ts' --include='*.tsx' | wc -l", 13, GREEN)
ap(tf, "Count type assertions:", 14, GRAY)
ap(tf, "grep -rn ' as [A-Z]' api/src web/src shared/src --include='*.ts' --include='*.tsx' | wc -l", 13, GREEN)
ap(tf, "Full type check:", 14, GRAY)
ap(tf, "pnpm type-check", 13, GREEN)

# ============================================================
# SLIDE 8: Cat 2 Before/After
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s)
tb(s, 0.5, 0.3, 12, 0.7, "Cat 2: Bundle Size \u2014 MET", 36, ACCENT, True)
tbl(s, 0.5, 1.2, 12, 4.5, [
    ["Metric", "Before", "After", "Change"],
    ["Largest chunk", "2,073 KB (94%)", "955 KB (32%)", "-54%"],
    ["Deferred", "0 KB", "784 KB", "New"],
    ["Route splits", "1 monolith", "23 chunks", "+22"],
    ["Emoji picker", "Always (271 KB)", "Lazy", "Deferred"],
    ["Highlight.js", "Always (195 KB)", "Lazy", "Deferred"],
], col_widths=[3, 3.5, 3, 2.5])

# ============================================================
# SLIDE 9: Cat 2 What Changed
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s)
tb(s, 0.5, 0.3, 12, 0.7, "Cat 2: What Changed", 36, ACCENT, True)
tf = tb(s, 0.5, 1.2, 12, 6.0, "Spec 2.2 \u2014 Route-Level Code Splitting", 16, GREEN, True)
ap(tf, "Converted all 23 page routes to React.lazy() with Suspense fallbacks. Before: visiting /login downloaded the entire app including editor, emoji picker, admin dashboard.", 14)
ap(tf, "Spec 2.3 \u2014 Lazy Emoji Picker", 16, GREEN, True)
ap(tf, "Deferred emoji-picker-react (271 KB). Only loaded when user clicks the emoji button.", 14)
ap(tf, "Spec 2.4 \u2014 Lazy Highlight.js", 16, GREEN, True)
ap(tf, "Deferred lowlight syntax highlighting (195 KB). Loaded only when a code block is present.", 14)
ap(tf, "Spec 2.5 \u2014 Lazy Upload Extensions", 16, GREEN, True)
ap(tf, "Deferred upload extensions (9 KB). Loaded on first file/image insert.", 14)
ap(tf, "Tradeoff: Total bundle +36% due to splitting overhead, but users download less per page.", 14, YELLOW)

# ============================================================
# SLIDE 10: Cat 2 How to Reproduce
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s)
tb(s, 0.5, 0.3, 12, 0.7, "Cat 2: How to Reproduce", 36, ACCENT, True)
tf = tb(s, 0.5, 1.2, 12, 5.0, "pnpm build:web", 16, GREEN)
ap(tf, "Vite prints all chunk sizes in the build output.", 14)
ap(tf, "Look for dist/assets/*.js files with sizes.", 14)
ap(tf, "Compare largest chunk and total JS size against baseline.", 14)

# ============================================================
# SLIDE 11: Cat 3 Before/After
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s)
tb(s, 0.5, 0.3, 12, 0.7, "Cat 3: API Response Time \u2014 MET", 36, ACCENT, True)
tbl(s, 0.5, 1.2, 12, 3.5, [
    ["Endpoint", "Before p99 (c=50)", "After p99 (c=50)", "Change"],
    ["Wiki docs", "142ms", "71ms", "-50%"],
    ["Issues (paginated)", "120ms", "76ms", "-37%"],
    ["Issues payload", "310 KB", "47 KB", "-85%"],
], col_widths=[3, 3, 3, 3])

# ============================================================
# SLIDE 12: Cat 3 What Changed
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s)
tb(s, 0.5, 0.3, 12, 0.7, "Cat 3: What Changed", 36, ACCENT, True)
tf = tb(s, 0.5, 1.2, 12, 6.0, "Spec 3.1 \u2014 Remove Content from Issues List", 16, GREEN, True)
ap(tf, "Dropped d.content from SELECT. List view only shows titles/metadata. Payload dropped 30%.", 14)
ap(tf, "Spec 3.2 \u2014 pg-pool Max Increase", 16, GREEN, True)
ap(tf, "Increased pool from 20 to 25 connections. At 50 concurrent requests, all 20 were occupied causing queuing. 5 extra connections eliminated contention.", 14)
ap(tf, "Spec 3.3 \u2014 Cursor-Based Pagination", 16, GREEN, True)
ap(tf, "Added cursor pagination to issues endpoint. Instead of all 218 issues (310 KB), returns 50 per page (47 KB).", 14)
ap(tf, "Tradeoff: Opening an issue requires separate content fetch. Pagination uses created_at sort.", 14, YELLOW)

# ============================================================
# SLIDE 13: Cat 3 How to Reproduce
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s)
tb(s, 0.5, 0.3, 12, 0.7, "Cat 3: How to Reproduce", 36, ACCENT, True)
tf = tb(s, 0.5, 1.2, 12, 5.5, "Prerequisites: pnpm db:seed, API running", 14, GRAY)
ap(tf, "autocannon -c 50 -d 30 -H 'Cookie: session=COOKIE' 'http://localhost:3001/api/documents?type=wiki'", 13, GREEN)
ap(tf, "autocannon -c 50 -d 30 -H 'Cookie: session=COOKIE' 'http://localhost:3001/api/issues?limit=50'", 13, GREEN)
ap(tf, "Replace COOKIE with valid session cookie from browser DevTools.", 14, GRAY)

# ============================================================
# SLIDE 14: Cat 4 Before/After
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s)
tb(s, 0.5, 0.3, 12, 0.7, "Cat 4: DB Queries \u2014 MET on ALL flows", 36, ACCENT, True)
tbl(s, 0.5, 1.2, 12, 5.0, [
    ["Flow", "Before", "After", "Change"],
    ["Main page", "25", "~16", "-36%"],
    ["View doc", "4", "~2", "-50%"],
    ["Issues", "5", "~4", "-20%"],
    ["Sprint board", "16", "~10", "-38%"],
    ["Search", "9", "~6", "-33%"],
    ["Auth/request", "3", "1", "-67%"],
], col_widths=[3, 3, 3, 3])

# ============================================================
# SLIDE 15: Cat 4 What Changed
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s)
tb(s, 0.5, 0.3, 12, 0.7, "Cat 4: What Changed", 36, ACCENT, True)
tf = tb(s, 0.5, 1.2, 12, 6.0, "Spec 4.1 \u2014 Auth Consolidation", 16, GREEN, True)
ap(tf, "Auth middleware ran 3 queries per request \u2014 60% of page load queries. Combined into single JOIN with 60s throttle on last_activity.", 14)
ap(tf, "Spec 4.2 \u2014 Remove Person JOIN", 16, GREEN, True)
ap(tf, "Issues list included expensive self-JOIN for archived assignee check. Removed from lists, kept in detail views.", 14)
ap(tf, "Spec 4.4 \u2014 Assignee Functional Index", 16, GREEN, True)
ap(tf, "Created B-tree index on (properties->>'assignee_id'). Dashboard queries went from sequential scan to index scan.", 14)
ap(tf, "Spec 4.5 \u2014 N+1 Batch Fix", 16, GREEN, True)
ap(tf, "Scope-changes had 1 query per removed issue. Replaced with single WHERE id = ANY($1) batch query.", 14)
ap(tf, "Tradeoff: last_activity 60s stale. Archived-assignee badge removed from lists.", 14, YELLOW)

# ============================================================
# SLIDE 16: Cat 4 How to Reproduce
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s)
tb(s, 0.5, 0.3, 12, 0.7, "Cat 4: How to Reproduce", 36, ACCENT, True)
tf = tb(s, 0.5, 1.2, 12, 5.0, "pnpm db:seed", 13, GREEN)
ap(tf, "Count queries using pg_stat_statements or application logging.", 14)
ap(tf, "EXPLAIN (ANALYZE, BUFFERS) on individual queries.", 14)
ap(tf, "Compare total query count per user flow against baseline.", 14)

# ============================================================
# SLIDE 17: Cat 5 Before/After
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s)
tb(s, 0.5, 0.3, 12, 0.7, "Cat 5: Test Coverage \u2014 EXCEEDED", 36, ACCENT, True)
tbl(s, 0.5, 1.2, 12, 4.5, [
    ["Metric", "Baseline", "Phase 2", "Phase 3"],
    ["Passing", "451 (30%)", "1,444 (99%)", "1,463 (99.5%)"],
    ["Web unit", "0/162", "138/151", "151/151"],
    ["E2E", "Not run", "855/855", "861/869"],
    ["Failures", "N/A", "13", "4 pre-existing"],
    ["Flaky", "N/A", "7", "4 pre-existing"],
    ["CI", "None", "None", "~5 min (8 shards)"],
], col_widths=[3, 3, 3, 3])
tb(s, 0.5, 5.8, 12, 0.6, "+1,012 tests restored. 30% to 99.5% operational.", 18, GREEN, True)

# ============================================================
# SLIDE 18: Cat 5 Phase 2 Fixes
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s)
tb(s, 0.5, 0.3, 12, 0.7, "Cat 5: Phase 2 \u2014 Infrastructure Fixes", 36, ACCENT, True)
tf = tb(s, 0.5, 1.2, 12, 6.0, "Spec 5.2 \u2014 Web Unit Tests (jsdom to happy-dom)", 16, GREEN, True)
ap(tf, "All 162 web tests broken by ERR_REQUIRE_ESM. html-encoding-sniffer dependency was ESM-only. Replaced jsdom environment with happy-dom. Result: 0/162 to 138/151 passing.", 14)
ap(tf, "Spec 5.3 \u2014 E2E Dynamic Import Fix", 16, GREEN, True)
ap(tf, "Playwright transforms static imports to require(). get-port is ESM-only. Changed to dynamic import() inside async function. Result: E2E suite unblocked, 855 tests confirmed passing.", 14)
ap(tf, "Commands:", 14, GRAY)
ap(tf, "pnpm test                        # API: 451 tests", 13, GREEN)
ap(tf, "pnpm --filter @ship/web test     # Web: 151 tests", 13, GREEN)

# ============================================================
# SLIDE 19: Cat 5 Phase 3 Fixes
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s)
tb(s, 0.5, 0.3, 12, 0.7, "Cat 5: Phase 3 \u2014 Test Quality Fixes", 36, ACCENT, True)
tf = tb(s, 0.5, 1.2, 12, 3.0, "13 Unit Test Failures Fixed", 16, GREEN, True)
ap(tf, "document-tabs.test.ts (9): stale assertions after tab rename. DetailsExtension.test.ts (3): missing child node imports. useSessionTimeout.test.ts (1): apiPost mock missing.", 14)
ap(tf, "10 of 14 Flaky E2E Tests Fixed", 16, GREEN, True)
ap(tf, "Root-caused and fixed with documented analysis. Common patterns:", 14)
tbl(s, 0.5, 4.2, 12, 3.0, [
    ["Pattern", "Tests", "Fix"],
    ["networkidle + WebSockets", "21+", "Remove"],
    ["waitForTimeout(N)", "15+", "toBeVisible/toPass"],
    ["Non-retrying assertions", "5", "Auto-retry"],
    ["Meta+ on Linux CI", "3", "ControlOrMeta+"],
    ["Silent-pass guards", "2", "Remove"],
], col_widths=[4.5, 3, 4.5])

# ============================================================
# SLIDE 20: Cat 5 CI Pipeline + Bugs
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s)
tb(s, 0.5, 0.3, 12, 0.7, "Cat 5: CI Pipeline + App Bugs Discovered", 36, ACCENT, True)
# Left side
tf = tb(s, 0.5, 1.2, 5.5, 5.0, "8-Shard CI Pipeline", 18, WHITE, True)
ap(tf, "Before: ~60 min local (1 worker, low RAM)", 14)
ap(tf, "After: ~5 min on GitHub Actions (8 shards x 4 workers)", 14)
ap(tf, "Each shard: own PostgreSQL + API + Vite", 14)
ap(tf, "Merged HTML report artifact", 14)
# Right side
tf = tb(s, 6.8, 1.2, 6.0, 5.0, "App Bugs Discovered", 18, RED, True)
ap(tf, "Yjs character truncation \u2014 last 3-4 chars lost during DB persistence", 14, RED)
ap(tf, "TipTap code block input rule \u2014 backtick rule doesn't fire for 2nd block", 14, RED)
ap(tf, "WCAG 1.4.13 vs 2.5.8 \u2014 focus controls trigger target-size violations", 14, RED)

# ============================================================
# SLIDE 21: Cat 5 How to Reproduce
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s)
tb(s, 0.5, 0.3, 12, 0.7, "Cat 5: How to Reproduce", 36, ACCENT, True)
tf = tb(s, 0.5, 1.2, 12, 5.5, "Unit tests (local):", 14, GRAY)
ap(tf, "pnpm test                        # API: 451 tests, ~12s", 13, GREEN)
ap(tf, "pnpm --filter @ship/web test     # Web: 151 tests, ~1.2s", 13, GREEN)
ap(tf, "E2E tests (CI):", 14, GRAY)
ap(tf, "git push origin cat5-test-coverage   # Triggers 8-shard pipeline", 13, GREEN)
ap(tf, "E2E tests (local):", 14, GRAY)
ap(tf, "pnpm test:e2e                    # 869 tests, requires Docker", 13, GREEN)

# ============================================================
# SLIDE 22: Cat 6 Before/After + What Changed
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s)
tb(s, 0.5, 0.3, 12, 0.7, "Cat 6: Error Handling \u2014 EXCEEDED (5/6)", 36, ACCENT, True)
tbl(s, 0.5, 1.2, 12, 4.5, [
    ["Failure", "Before", "After"],
    ["No ErrorBoundary", "Blank screen", "Recovery UI + retry"],
    ["WS storm on 429", "3s retry forever", "Exponential 3s-60s"],
    ["Rate-limit leak", "Never decremented", "release() on close"],
    ["Title >255", "Silently reverts", "maxLength + counter"],
    ["Save ROLLBACK", ".catch(() => {})", "Retry 3x + toast"],
], col_widths=[3.5, 4, 4.5])
tb(s, 0.5, 6.0, 12, 0.5, "Phase 3: Spec 6.6 \u2014 Real-time title sync + useAutoSave throttle fix", 14, GRAY)

# ============================================================
# SLIDE 23: Cat 6 How to Reproduce
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s)
tb(s, 0.5, 0.3, 12, 0.7, "Cat 6: How to Reproduce", 36, ACCENT, True)
tf = tb(s, 0.5, 1.2, 12, 5.5, "1. ErrorBoundary: Temporarily throw in a provider, verify recovery UI", 14)
ap(tf, "2. WS backoff: Open DevTools Network, throttle to offline, verify timing", 14)
ap(tf, "3. Title guard: Type 256+ characters, verify counter and enforcement", 14)
ap(tf, "4. Save failure: Disconnect DB mid-save, verify toast appears", 14)
ap(tf, "5. Title sync: Open same doc in 2 tabs, change title, verify sync", 14)

# ============================================================
# SLIDE 24: Cat 7 Before/After + What Changed
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s)
tb(s, 0.5, 0.3, 12, 0.7, "Cat 7: Accessibility \u2014 MET (13/13)", 36, ACCENT, True)
tbl(s, 0.3, 1.2, 12.7, 5.5, [
    ["Fix", "Page", "What Changed"],
    ["aria-expanded", "/documents/:id", "Removed invalid attr"],
    ["Contrast (4 fixes)", "Multiple", "~1.6:1 to 4.5:1"],
    ["4 aria-labels", "/documents/:id", "Search, emoji, backlog, chips"],
    ["Label association", "/documents/:id", "useId() for inputs"],
    ["TabBar keyboard", "/my-week", "Arrow/Home/End + roving tabindex"],
    ["20 decorative SVGs", "/projects", "aria-hidden on icons"],
    ["Focus ring", "/projects", "focus-visible on selectable rows"],
], col_widths=[3.5, 3.5, 5.7])

# ============================================================
# SLIDE 25: Cat 7 WCAG Conflict
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s)
tb(s, 0.5, 0.3, 12, 0.7, "Cat 7: WCAG Conflict (1.4.13 vs 2.5.8)", 36, ACCENT, True)
tf = tb(s, 0.5, 1.2, 12, 6.0, "WCAG 1.4.13: Hover controls must also show on keyboard focus", 16, WHITE)
ap(tf, "Requires tabIndex on tree items so they can receive focus", 14, GRAY)
ap(tf, "", 14)
ap(tf, "WCAG 2.5.8: All focusable elements must be at least 24x24px", 16, WHITE)
ap(tf, "Tree items: <li> at 6px, links at 20px, buttons at 16-18px", 14, RED)
ap(tf, "", 14)
ap(tf, "Attempted fix (F3.11): Added tabIndex. Axe-core immediately flagged target-size.", 14, YELLOW)
ap(tf, "Reverted (F3.30): Documented as future work.", 14, YELLOW)
ap(tf, "Correct fix: Increase all touch targets to 24px FIRST, then add tabIndex.", 14, GREEN)

# ============================================================
# SLIDE 26: Cat 7 How to Reproduce
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s)
tb(s, 0.5, 0.3, 12, 0.7, "Cat 7: How to Reproduce", 36, ACCENT, True)
tf = tb(s, 0.5, 1.2, 12, 5.5, "Axe-core E2E tests:", 14, GRAY)
ap(tf, "pnpm test:e2e --grep 'axe-core'", 13, GREEN)
ap(tf, "Lighthouse (Chrome DevTools):", 14, GRAY)
ap(tf, "1. Open DevTools > Lighthouse tab", 14)
ap(tf, "2. Select Accessibility only, Desktop viewport", 14)
ap(tf, "3. Run on /documents/:id, /my-week, /projects", 14)

# ============================================================
# SLIDE 27: Discovery 1
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s)
tb(s, 0.5, 0.3, 12, 0.7, "Discovery: Everything is a Document", 36, ACCENT, True)
tf = tb(s, 0.5, 1.2, 12, 6.0, "Found in: shared/src/types/document.ts (lines 34-330)", 14, GRAY)
ap(tf, "ShipShape stores every entity type \u2014 wikis, issues, projects, programs, sprints, weeks \u2014 in a single documents table with a document_type discriminator. Any document can link to any other. The TypeScript discriminated union narrows types automatically on switch(doc.type).", 14)
ap(tf, "", 14)
ap(tf, "Application: openEMR", 18, GREEN, True)
ap(tf, "Patient visits, prescriptions, lab results, referrals could all be documents with type-specific JSON schemas. Lab results from external systems stored directly and linked to encounters. Doctors see the full picture without switching modules.", 14)

# ============================================================
# SLIDE 28: Discovery 2 + 3
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s)
tb(s, 0.5, 0.3, 12, 0.7, "Discoveries: Playwright + useAutoSave", 36, ACCENT, True)
# Left side
tf = tb(s, 0.5, 1.2, 5.8, 5.5, "Playwright E2E Framework", 18, GREEN, True)
ap(tf, "Found in: e2e/ (869 tests)", 13, GRAY)
ap(tf, "Auto-retrying assertions vs one-shot reads", 14)
ap(tf, "Per-worker isolation via testcontainers", 14)
ap(tf, "8-shard CI: 60 min to 5 min", 14)
ap(tf, "Powerful but unforgiving on timing", 14)
# Right side
tf = tb(s, 6.8, 1.2, 6.0, 5.5, "useAutoSave Hook", 18, GREEN, True)
ap(tf, "Found in: web/src/hooks/useAutoSave.ts", 13, GRAY)
ap(tf, "Throttle: saves every 500ms while typing", 14)
ap(tf, "Queue: changes during in-flight saves", 14)
ap(tf, "Sequence #: prevents stale overwrites", 14)
ap(tf, "Retry: exponential backoff on failure", 14)

# ============================================================
# SLIDE 29: AI Cost Analysis
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s)
tb(s, 0.5, 0.3, 12, 0.7, "AI Cost Analysis", 36, ACCENT, True)
tbl(s, 0.5, 1.2, 12, 3.5, [
    ["Item", "Details"],
    ["Primary tool", "Claude Code MAX ($200/mo)"],
    ["Secondary", "ChatGPT"],
    ["Estimated tokens", "~10M+ across all sessions"],
    ["Project time", "~70 hours"],
], col_widths=[4, 8])
tb(s, 0.5, 4.8, 12, 0.5, "95% AI-generated code, 5% hand-written", 16, WHITE, True)
tb(s, 0.5, 5.4, 12, 0.5, "Human role: architect and quality gate. AI role: implementer.", 14, GRAY)

# ============================================================
# SLIDE 30: AI Reflection
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s)
tb(s, 0.5, 0.3, 12, 0.7, "AI Reflection", 36, ACCENT, True)
tf = tb(s, 0.5, 1.2, 12, 6.0, "Most helpful: Planning, automation, documentation, CI sharding", 14, GREEN)
ap(tf, "Least helpful: Code changes without specs \u2014 fabricated results, took shortcuts", 14, RED)
ap(tf, "Key lesson: AI amplifies direction. Good workflows accelerate; bad ones too.", 14, YELLOW)
ap(tf, "", 14)
ap(tf, "Overrides required:", 16, WHITE, True)
ap(tf, "Commit discipline \u2014 force-merged without structured messages", 14)
ap(tf, "Fabricated test results \u2014 solved by pushing to CI", 14)
ap(tf, "Regression-causing fixes \u2014 3 reverted after CI verification", 14)
ap(tf, "Unauthorized pushes \u2014 had to require approval for every commit", 14)
ap(tf, "Over-engineering \u2014 CLAUDE.md rule: don't add unnecessary code", 14)

# ============================================================
# SLIDE 31: Remaining + Future
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s)
tb(s, 0.5, 0.3, 12, 0.7, "Remaining Issues + Future Work", 36, ACCENT, True)
tb(s, 0.5, 1.1, 12, 0.4, "4 Hard Failures (pre-existing):", 16, WHITE, True)
tbl(s, 0.5, 1.6, 12, 2.8, [
    ["Test", "Area", "Root Cause"],
    ["a11y:407", "Accessibility", "WCAG conflict"],
    ["syntax:154", "Code Blocks", "Yjs truncation"],
    ["data:356", "Mentions", "Yjs persistence"],
    ["drag:310", "Drag Handle", "Element timing"],
], col_widths=[3, 4, 5])
tb(s, 0.5, 4.5, 12, 0.4, "Future Work:", 16, WHITE, True)
tbl(s, 0.5, 5.0, 12, 2.3, [
    ["Item", "Description"],
    ["Yjs persistence", "Fix character truncation"],
    ["WCAG targets", "24px minimums then tabIndex"],
    ["Multi-user tests", "2 users editing same doc"],
    ["CAIA/OAuth", "Government login zero coverage"],
    ["API coverage", "33% to 60%+ branches"],
], col_widths=[4, 8])

# ============================================================
# SLIDE 32: Thank You
# ============================================================
s = prs.slides.add_slide(BLANK_LAYOUT)
set_bg(s)
tb(s, 0.5, 1.5, 12.333, 1.0, "ShipShape", 48, ACCENT, True, PP_ALIGN.CENTER)
tb(s, 0.5, 2.7, 12.333, 0.8, "7 categories. 24 specs. All targets met.", 24, GREEN, False, PP_ALIGN.CENTER)
tb(s, 0.5, 3.7, 12.333, 0.5, "Final Report: docs/final-audit-report.md", 18, GRAY, False, PP_ALIGN.CENTER)
tb(s, 0.5, 4.3, 12.333, 0.5, "GitHub: github.com/fsyeddev-a11y/ShipShape", 18, GRAY, False, PP_ALIGN.CENTER)
tb(s, 0.5, 4.9, 12.333, 0.5, "Deployed: shipshape-prod-web.onrender.com", 18, GRAY, False, PP_ALIGN.CENTER)

# ============================================================
# Save
# ============================================================
output_path = os.path.join(os.path.dirname(__file__), "ShipShape-Final-Presentation.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
